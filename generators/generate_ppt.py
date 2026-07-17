from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import os
from datetime import datetime

def add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = data.get('title', 'QA Test Report')
    subtitle.text = f"Project: {data.get('project', 'N/A')}\nTester: {data.get('tester', 'N/A')}\nDate: {data.get('date', 'N/A')}"

def add_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"

    content = slide.placeholders[1]
    summary = data.get('summary', {})
    total = summary.get('total', 0)
    passed = summary.get('passed', 0)
    failed = summary.get('failed', 0)
    pass_rate = (passed / total * 100) if total > 0 else 0

    text_frame = content.text_frame
    text_frame.text = f"Total Tests: {total}\n"
    text_frame.text += f"\u2705 Passed: {passed}\n"
    text_frame.text += f"\u274c Failed: {failed}\n"
    text_frame.text += f"\U0001f4ca Pass Rate: {pass_rate:.1f}%\n"

def add_tests_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Test Cases"

    content = slide.placeholders[1]
    text_frame = content.text_frame

    tests = data.get('tests', [])
    for test in tests:
        status = test.get('status', '')
        emoji = "\u2705" if status.lower() == 'pass' else "\u274c"
        text_frame.text += f"{emoji} [{test.get('id', '')}] {test.get('scenario', '')} - {status}\n"

def generate_ppt(data, output_path):
    prs = Presentation()

    add_title_slide(prs, data)
    add_summary_slide(prs, data)
    add_tests_slide(prs, data)

    prs.save(output_path)

def generate_ppt_from_file(json_file='data.json', output_dir='output'):
    try:
        with open(json_file, 'r') as f:
            data = json.load(f)

        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, 'report.pptx')
        generate_ppt(data, output_path)
        print(f"\u2705 PPTX generated: {output_path}")
        return output_path
    except Exception as e:
        print(f"\u274c Error generating PPTX: {e}")
        return None

if __name__ == "__main__":
    generate_ppt_from_file()
