from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from datetime import datetime


def generate_text_pdf(text, output_path, title="Document"):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, title, 0, 1, "C")
    pdf.set_font("Helvetica", "I", 10)
    pdf.cell(0, 6, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", 0, 1, "C")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 10)

    for line in text.split("\n"):
        if not line.strip():
            pdf.ln(4)
            continue

        pdf.set_x(10)
        safe = line.encode("latin-1", "replace").decode("latin-1")
        w = pdf.w - pdf.l_margin - pdf.r_margin

        try:
            if pdf.get_string_width(safe) > w:
                pdf.multi_cell(w=w, h=5, text=safe)
            else:
                pdf.cell(w, 5, safe, 0, 1)
        except Exception:
            pdf.set_font("Helvetica", "", 8)
            try:
                pdf.multi_cell(w=w, h=4, text=safe[:500])
            except Exception:
                pass
            pdf.set_font("Helvetica", "", 10)

    pdf.output(output_path)


def generate_text_pptx(text, output_path, title="Document"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, title)

    lines = text.split("\n")
    chunk_size = 20
    chunks = [lines[i:i + chunk_size] for i in range(0, len(lines), chunk_size)]

    for i, chunk in enumerate(chunks):
        if i > 0:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _set_title(slide, f"{title} (cont.)")

        body = "\n".join(chunk)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, line in enumerate(chunk):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    prs.save(output_path)


def generate_text_docx(text, output_path, title="Document"):
    doc = Document()
    h = doc.add_heading(title, 0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    meta = doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    for line in text.split("\n"):
        if line.strip():
            p = doc.add_paragraph(line)
            for run in p.runs:
                run.font.size = DocxPt(11)
        else:
            doc.add_paragraph()

    doc.save(output_path)


def _set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0x2F, 0x54, 0x96)
